"""
RAG pipeline for Professor Tux — lecture slide ingestion and retrieval.
"""

import os
import uuid
import hashlib
import logging
from pathlib import Path
from typing import Optional
from datetime import datetime, timezone
from dataclasses import dataclass, asdict

import chromadb
from chromadb.config import Settings as ChromaSettings
from sentence_transformers import SentenceTransformer

logger = logging.getLogger("professor_tux.rag")

CHROMA_PERSIST_DIR = os.getenv("CHROMA_PERSIST_DIR", "./data/chromadb")
EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", "all-MiniLM-L6-v2")
CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "500"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "100"))
TOP_K_RESULTS = int(os.getenv("TOP_K_RESULTS", "5"))
UPLOAD_DIR = os.getenv("UPLOAD_DIR", "./data/uploads")


@dataclass
class LectureDocument:
    doc_id: str
    filename: str
    course: Optional[str] = None
    lecture_title: Optional[str] = None
    file_type: str = ""
    num_chunks: int = 0
    num_pages: int = 0
    uploaded_at: str = ""
    file_hash: str = ""
    def to_dict(self) -> dict: return asdict(self)


@dataclass
class RetrievedContext:
    text: str
    source_filename: str
    lecture_title: Optional[str]
    course: Optional[str]
    page_or_slide: Optional[int]
    relevance_score: float


# ── Extractors ───────────────────────────────────────────────────────

def extract_text_from_pdf(file_path: str) -> list[dict]:
    import fitz
    pages = []
    doc = fitz.open(file_path)
    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if text: pages.append({"page": i, "text": text})
    doc.close()
    return pages


def extract_text_from_pptx(file_path: str) -> list[dict]:
    from pptx import Presentation
    slides = []
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        if slide.shapes.title and slide.shapes.title.text.strip():
            texts.append(f"[Slide Title: {slide.shapes.title.text.strip()}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t: texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    rt = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if rt: texts.append(rt)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes: texts.append(f"[Speaker Notes: {notes}]")
        combined = "\n".join(texts)
        if combined.strip(): slides.append({"page": i, "text": combined})
    return slides


def extract_text_from_txt(file_path: str) -> list[dict]:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read().strip()
    return [{"page": 1, "text": content}] if content else []


EXTRACTORS = {
    ".pdf": extract_text_from_pdf,
    ".pptx": extract_text_from_pptx,
    ".ppt": extract_text_from_pptx,
    ".txt": extract_text_from_txt,
    ".md": extract_text_from_txt,
}


# ── Chunking ─────────────────────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    if len(text) <= chunk_size: return [text]
    chunks, start = [], 0
    while start < len(text):
        end = start + chunk_size
        if end < len(text):
            for bc in [". ", ".\n", "? ", "!\n", "\n\n", "\n"]:
                b = text.rfind(bc, start + chunk_size // 2, end)
                if b != -1: end = b + len(bc); break
        chunk = text[start:end].strip()
        if chunk: chunks.append(chunk)
        start = end - overlap
        if start >= len(text): break
    return chunks


# ── Knowledge Base ───────────────────────────────────────────────────

class LectureKnowledgeBase:
    def __init__(self):
        self._embedder: Optional[SentenceTransformer] = None
        self._chroma_client = None
        self._collection = None
        self._documents: dict[str, LectureDocument] = {}

    @property
    def is_loaded(self) -> bool:
        return self._embedder is not None and self._collection is not None

    def initialize(self):
        os.makedirs(CHROMA_PERSIST_DIR, exist_ok=True)
        os.makedirs(UPLOAD_DIR, exist_ok=True)
        logger.info("⏳ Loading embedding model: %s", EMBEDDING_MODEL)
        self._embedder = SentenceTransformer(EMBEDDING_MODEL)
        logger.info("✅ Embedding model loaded (dim=%d)", self._embedder.get_sentence_embedding_dimension())
        self._chroma_client = chromadb.PersistentClient(
            path=CHROMA_PERSIST_DIR, settings=ChromaSettings(anonymized_telemetry=False))
        self._collection = self._chroma_client.get_or_create_collection(
            name="lecture_slides", metadata={"hnsw:space": "cosine"})
        logger.info("✅ ChromaDB ready (%d chunks)", self._collection.count())
        self._rebuild_doc_index()

    def _rebuild_doc_index(self):
        if self._collection.count() == 0: return
        results = self._collection.get(include=["metadatas"])
        seen = {}
        for meta in results["metadatas"]:
            did = meta.get("doc_id")
            if did and did not in seen:
                seen[did] = LectureDocument(
                    doc_id=did, filename=meta.get("filename", "?"),
                    course=meta.get("course"), lecture_title=meta.get("lecture_title"),
                    file_type=meta.get("file_type", ""), uploaded_at=meta.get("uploaded_at", ""),
                    file_hash=meta.get("file_hash", ""))
        self._documents = seen

    def ingest_file(self, file_path: str, course: Optional[str] = None,
                    lecture_title: Optional[str] = None) -> LectureDocument:
        if not self.is_loaded: raise RuntimeError("KB not initialized")
        fp = Path(file_path)
        suffix = fp.suffix.lower()
        if suffix not in EXTRACTORS:
            raise ValueError(f"Unsupported: {suffix}")

        file_hash = hashlib.sha256(fp.read_bytes()).hexdigest()[:16]
        for doc in self._documents.values():
            if doc.file_hash == file_hash: return doc

        pages = EXTRACTORS[suffix](str(fp))
        if not pages: raise ValueError(f"No text from {fp.name}")

        doc_id = str(uuid.uuid4())
        now = datetime.now(timezone.utc).isoformat()
        all_chunks, all_meta, all_ids = [], [], []

        for pd in pages:
            for ci, chunk in enumerate(chunk_text(pd["text"])):
                cid = f"{doc_id}_p{pd['page']}_c{ci}"
                all_chunks.append(chunk)
                all_ids.append(cid)
                all_meta.append({
                    "doc_id": doc_id, "filename": fp.name, "course": course or "",
                    "lecture_title": lecture_title or "", "page_or_slide": pd["page"],
                    "chunk_index": ci, "file_type": suffix, "file_hash": file_hash,
                    "uploaded_at": now,
                })

        embeddings = self._embedder.encode(all_chunks, show_progress_bar=False).tolist()
        self._collection.add(ids=all_ids, documents=all_chunks,
                             embeddings=embeddings, metadatas=all_meta)

        doc = LectureDocument(doc_id=doc_id, filename=fp.name, course=course,
                              lecture_title=lecture_title, file_type=suffix,
                              num_chunks=len(all_chunks), num_pages=len(pages),
                              uploaded_at=now, file_hash=file_hash)
        self._documents[doc_id] = doc
        logger.info("✅ Ingested '%s': %d pages → %d chunks", fp.name, len(pages), len(all_chunks))
        return doc

    def search(self, query: str, top_k: int = TOP_K_RESULTS,
               course_filter: Optional[str] = None, doc_id_filter: Optional[str] = None) -> list[RetrievedContext]:
        if not self.is_loaded or self._collection.count() == 0: return []
        qe = self._embedder.encode([query]).tolist()
        wf = None
        conds = []
        if course_filter: conds.append({"course": {"$eq": course_filter}})
        if doc_id_filter: conds.append({"doc_id": {"$eq": doc_id_filter}})
        if len(conds) == 1: wf = conds[0]
        elif len(conds) > 1: wf = {"$and": conds}

        results = self._collection.query(
            query_embeddings=qe, n_results=min(top_k, self._collection.count()),
            where=wf, include=["documents", "metadatas", "distances"])

        contexts = []
        for i in range(len(results["ids"][0])):
            meta = results["metadatas"][0][i]
            sim = 1 - (results["distances"][0][i] / 2)
            contexts.append(RetrievedContext(
                text=results["documents"][0][i], source_filename=meta.get("filename", "?"),
                lecture_title=meta.get("lecture_title") or None, course=meta.get("course") or None,
                page_or_slide=meta.get("page_or_slide"), relevance_score=round(sim, 4)))
        return contexts

    def format_context_for_prompt(self, contexts: list[RetrievedContext], min_rel: float = 0.3) -> str:
        relevant = [c for c in contexts if c.relevance_score >= min_rel]
        if not relevant: return ""
        lines = ["━━━ RELEVANT LECTURE MATERIAL ━━━",
                 "Use this material to ground your teaching in what has been covered in class.\n"]
        for c in relevant:
            src = f"📎 {c.source_filename}"
            if c.lecture_title: src += f" — {c.lecture_title}"
            if c.page_or_slide: src += f" (Slide {c.page_or_slide})"
            src += f" [{c.relevance_score:.0%}]"
            lines.extend([src, c.text, ""])
        lines.append("━━━ END LECTURE MATERIAL ━━━")
        return "\n".join(lines)

    def list_documents(self) -> list[dict]:
        return [d.to_dict() for d in self._documents.values()]

    def get_document(self, doc_id: str) -> Optional[dict]:
        d = self._documents.get(doc_id)
        return d.to_dict() if d else None

    def delete_document(self, doc_id: str) -> bool:
        if doc_id not in self._documents: return False
        results = self._collection.get(where={"doc_id": {"$eq": doc_id}}, include=[])
        if results["ids"]: self._collection.delete(ids=results["ids"])
        del self._documents[doc_id]
        return True

    def get_stats(self) -> dict:
        return {"total_documents": len(self._documents),
                "total_chunks": self._collection.count() if self._collection else 0,
                "embedding_model": EMBEDDING_MODEL, "chunk_size": CHUNK_SIZE,
                "chunk_overlap": CHUNK_OVERLAP}
